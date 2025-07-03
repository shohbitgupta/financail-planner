from pptx import Presentation
import re
import os

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# ✅ Step 1: Read slide content from a file
file_path = os.path.join(os.getcwd(), "Financial_Planner_AI_Agent_Presentation.md")  # Ensure this file is in your root directory
with open(file_path, "r", encoding="utf-8") as file:
    presentation_text = file.read()

slides_content = re.split(r'\n## Slide \d+: ', presentation_text)
slides_content = [slide.strip() for slide in slides_content if slide.strip()]

for i, content in enumerate(slides_content):
    lines = content.strip().split('\n')
    title = lines[0].replace('### ', '').replace('🏗️ ', '').strip()
    body_lines = lines[1:]

    if i == 0:
        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        subtitle = slide.placeholders[1]
        title_placeholder.text = title
        subtitle.text = "\n".join(body_lines)
    else:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title
        tf = body_shape.text_frame
        tf.clear()

        for line in body_lines:
            if line.startswith('- ') or line.startswith('• '):
                p = tf.add_paragraph()
                p.text = line[2:]
                p.level = 0
            elif line.startswith('  - '):
                p = tf.add_paragraph()
                p.text = line.strip()[4:]
                p.level = 1
            elif line.strip():
                p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 0

output_path = os.path.join(os.getcwd(), "Financial_Planner_AI_Agent_Architecture.pptx")
prs.save(output_path)
print(f"Presentation saved as: {output_path}")
